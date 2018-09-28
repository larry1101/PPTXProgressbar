from pptx.dml.color import RGBColor
from pptx.util import Inches, Length
from pptx.enum.shapes import MSO_SHAPE


def add_rect(slide, left, top, width, height, fill='89ABBF'):
    # left, top, width, height = Inches(0), Inches(0), Inches(13.3326), Inches(0.10625)  # 预设位置及大小
    # left, top, width, height = Inches(0), Inches(0), Length(12192000), Length(100000)  # 预设位置及大小
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)  # 在指定位置按预设值添加类型为PENTAGON的形状

    shape_fill = shape.fill
    shape_fill.solid()
    # shape_fill.fore_color.rgb = RGBColor(180, 216, 222)
    shape_fill.fore_color.rgb = RGBColor.from_string(fill)

    shape_border_line = shape.line
    shape_border_line.fill.background()


def add_text(slide, left, top, width, height, size, text=''):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)  # 在指定位置按预设值添加类型为PENTAGON的形状
    shape_fill = shape.fill
    shape_fill.background()
    shape_border_line = shape.line
    shape_border_line.fill.background()

    shape.text = text
    shape_text_font = shape.text_frame.paragraphs[0].font
    shape_text_font.size = size
